# app.py
# ==============================================
# Intelligent Document Assistant
# Chat history memory · Separate models per task
# Azure AD auth · Session rename · Display names
# ==============================================

import base64
# import configparser
from dotenv import load_dotenv
import datetime as dt
import io
import json
import logging
import os
import uuid
from concurrent.futures import ThreadPoolExecutor, as_completed

import csv
from docx import Document as DocxDocument      # Word (.docx) support
import openpyxl                                 # Excel (.xlsx) support
from pptx import Presentation                   # PowerPoint (.pptx) support
from pypdf import PdfReader, PdfWriter


import faiss
import msal
import numpy as np
from azure.ai.documentintelligence import DocumentIntelligenceClient
from azure.ai.documentintelligence.models import DocumentContentFormat
from azure.core.credentials import AzureKeyCredential
from azure.core.exceptions import HttpResponseError
from azure.core.exceptions import ResourceNotFoundError
from azure.storage.blob import BlobServiceClient
from flask import (
    Flask,
    Response,
    jsonify,
    redirect,
    render_template,
    request,
    session,
    url_for,
    stream_with_context,
)
from flask_session import Session
from openai import AzureOpenAI, RateLimitError

load_dotenv()   # pulls values from .env into os.environ

# ------------------------------------------------
# 0.  Logging - start
# ------------------------------------------------
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)

# ------------------------------------------------
# 1.  Flask & Session
# ------------------------------------------------
app = Flask(__name__)

# Session backend:
# - Default: filesystem (works if Azure keeps you on a single instance)
# - If FLASK_SESSION_TYPE is set (e.g. "redis"), Flask-Session will use that.
app.config["SESSION_TYPE"] = os.getenv("FLASK_SESSION_TYPE", "filesystem")

# IMPORTANT:
# Use a stable secret key from environment so all instances can
# decrypt the same session cookies.

secret_from_env = os.getenv("FLASK_SECRET_KEY")
if not secret_from_env:
    raise RuntimeError(
        "Missing environment variable: FLASK_SECRET_KEY. "
        "Set this in your local .env and in Azure App Settings so "
        "sessions/cookies remain stable across instances."
    )
app.config["SECRET_KEY"] = secret_from_env

# Global upload + session safety limits for prod:
# - MAX_CONTENT_LENGTH_MB: hard cap per HTTP request (Flask rejects earlier).
# - SESSION_LIFETIME_HOURS: how long auth/session cookies stay valid.
max_mb = int(os.getenv("MAX_CONTENT_LENGTH_MB", "4096"))  # default 4096 MB per request
app.config["MAX_CONTENT_LENGTH"] = max_mb * 1024 * 1024

app.config["SESSION_PERMANENT"] = False
app.config["PERMANENT_SESSION_LIFETIME"] = dt.timedelta(
    hours=int(os.getenv("SESSION_LIFETIME_HOURS", "8"))
)

# Cookie behaviour – keep it robust behind proxies / Azure portal iframe
# If you ONLY ever run top-level (not in iframe), you can change SameSite to 'Lax'.
app.config["SESSION_COOKIE_SECURE"] = True          # cookies only over HTTPS
app.config["SESSION_COOKIE_SAMESITE"] = "None"      # allow cookies in iframe contexts

Session(app)

# ------------------------------------------------
# 2.  Config (env variables required; unchanged)
# ------------------------------------------------
def need(name: str) -> str:
    """Grab required env var or raise a clear error."""
    val = os.getenv(name)
    if not val:
        raise RuntimeError(f"Missing environment variable: {name}")
    return val

# Azure AI
AI_FOUNDRY_KEY      = need("AI_FOUNDRY_KEY")
AI_FOUNDRY_ENDPOINT = need("AI_FOUNDRY_ENDPOINT").rstrip("/")
SUMMARY_DEPLOYMENT  = need("SUMMARY_DEPLOYMENT_NAME")
CHAT_DEPLOYMENT     = need("CHAT_DEPLOYMENT_NAME")
EMBED_DEPLOYMENT    = need("EMBED_DEPLOYMENT_NAME")
API_VER             = "2024-02-01"
CHAT_HISTORY_TURNS  = int(os.getenv("CHAT_HISTORY_TURNS", "6"))
# Document Intelligence
DOCUMENT_INTELLIGENCE_ENDPOINT = need("DOCUMENT_INTELLIGENCE_ENDPOINT")
DOCUMENT_INTELLIGENCE_KEY      = need("DOCUMENT_INTELLIGENCE_KEY")

# Blob Storage
BLOB_CONNECTION_STRING = need("BLOB_CONNECTION_STRING")
BLOB_CONTAINER_NAME    = need("BLOB_CONTAINER_NAME")

# App base URL – single source of truth for redirects.
# Example: "https://ida.company.com" or "https://yourapp.azurewebsites.net"
APP_BASE_URL = need("APP_BASE_URL").rstrip("/")

# Azure AD
AAD_CLIENT_ID     = need("AAD_CLIENT_ID")
AAD_CLIENT_SECRET = need("AAD_CLIENT_SECRET")
AAD_AUTHORITY     = need("AAD_AUTHORITY").rstrip("/")
AAD_REDIRECT      = "/get_token"
AAD_SCOPE         = ["User.Read"]

# Fixed redirect URI that MUST match the Entra app registration exactly.
AAD_REDIRECT_URI = f"{APP_BASE_URL}{AAD_REDIRECT}"

# ------------------------------------------------
# 3.  Clients
# ------------------------------------------------
llm_client = AzureOpenAI(
    api_key=AI_FOUNDRY_KEY,
    api_version=API_VER,
    azure_endpoint=AI_FOUNDRY_ENDPOINT,
    max_retries=5,
)
doc_client = DocumentIntelligenceClient(
    endpoint=DOCUMENT_INTELLIGENCE_ENDPOINT, credential=AzureKeyCredential(DOCUMENT_INTELLIGENCE_KEY)
)
blob_service = BlobServiceClient.from_connection_string(BLOB_CONNECTION_STRING)
container = blob_service.get_container_client(BLOB_CONTAINER_NAME)
if not container.exists():
    container.create_container()

msal_client = msal.ConfidentialClientApplication(
    AAD_CLIENT_ID, authority=AAD_AUTHORITY, client_credential=AAD_CLIENT_SECRET
)

# ------------------------------------------------
# 4.  Prompts
# ------------------------------------------------
SYSTEM_PROMPT_SUMMARY = (
    "You are a document summarization assistant. "
    "Produce a clear, detailed, and professional summary strictly based on the provided text. "
    "Write for a general business audience reviewing uploaded materials.\n\n"

    "ROLE & ACCURACY RULES:\n"
    "- Identify the primary people, organizations, and records described in the documents.\n"
    "- Distinguish between document subjects and administrative contacts when that difference is clear.\n"
    "- Do NOT invent names, dates, numbers, diagnoses, or events.\n\n"

    "CONTEXTUAL INFERENCE GUIDANCE (IMPORTANT):\n"
    "- You MAY draw reasonable conclusions that are clearly supported by the document context.\n"
    "- If information is strongly implied across multiple documents, summarize it confidently.\n"
    "- If information is unclear, conflicting, or only partially referenced, state that explicitly.\n"
    "- Do NOT downgrade clear facts to 'Not stated' merely because wording varies across documents.\n\n"

    "WRITING GUIDELINES:\n"
    "- Do NOT reference the summarization process or uncertainty mechanics.\n"
    "- Be specific and detailed when the documents support it.\n"
    "- Be concise but thorough so the summary is practical and easy to review.\n\n"

    "OUTPUT FORMAT (use bold section titles and bullet points):\n"
    "**Document Types**\n"
    "- Describe the kinds of documents present (for example correspondence, reports, invoices, records, or forms).\n\n"
    "**People and Organizations**\n"
    "- List named individuals, teams, vendors, clients, or organizations and note their roles when stated.\n\n"
    "**Key Facts**\n"
    "- Summarize the main events, requests, findings, or decisions described in the documents.\n\n"
    "**Important Dates and Timeline**\n"
    "- List significant dates and the related events in chronological order when possible.\n\n"
    "**Amounts, Metrics, or Commitments**\n"
    "- Summarize any notable amounts, quantities, deadlines, obligations, or tracked metrics that are documented.\n\n"
    "**Open Questions or Missing Information**\n"
    "- Identify genuinely missing or unclear information that would help a reviewer understand the document set."
)

SYSTEM_PROMPT_CHAT = (
    "You are a document analysis assistant. "
    "Your role is to help users locate, verify, and understand specific information "
    "within the provided documents.\n\n"

    "GENERAL RULES:\n"
    "- Use ONLY the provided document context and prior conversation.\n"
    "- Do NOT invent facts or assume information that is not supported by the documents.\n\n"

    "QUESTION HANDLING:\n"
    "- Carefully identify what the user is asking (dates, counts, people, obligations, findings, or history).\n"
    "- Review the document context for explicit statements and clearly supported implications.\n"
    "- For time-based questions (e.g., 'before a certain date'), evaluate only events documented prior to that date.\n"
    "- For counting questions, count only explicitly documented events.\n\n"

    "ABSENCE & HISTORY CHECKS:\n"
    "- If asked whether something exists in the record:\n"
    "  • Answer 'Yes' only if explicitly documented.\n"
    "  • Answer 'No documentation found' if the documents are silent.\n"
    "  • Do NOT assume absence means none existed.\n\n"

    "EVIDENCE:\n"
    "- Support answers using brief references to the document context "
    "(source file and page number when available).\n"
    "- If no supporting text exists, state that the information is not documented.\n\n"

    "TONE:\n"
    "- Be clear, direct, and professional.\n"
    "- Keep answers useful for a general document review workflow."
)


# ------------------------------------------------
# 5.  Helpers – identity, blobs, sessions
# ------------------------------------------------
def _sanitize(s: str) -> str:
    """
    Make a string safe for use in blob paths:
    - allow alphanumeric, dash, underscore, dot
    - everything else is stripped
    """
    if not isinstance(s, str):
        s = str(s or "")
    return "".join(c for c in s if c.isalnum() or c in "-_.")


def current_claims() -> dict | None:
    c = session.get("user")
    return c if isinstance(c, dict) else None


def current_user_id() -> str | None:
    """
    Stable ID *for storage & gating*.

    New behavior:
      - Prefer human-readable identity for the user folder in Blob:
          preferred_username / upn / email / unique_name / name
      - Fallback to the old tid:oid combo if none of the above exist.

    This means blob layout will look like:
      user-sessions/<email-or-name>/<session_id>/...

    while still being stable per user.
    """
    claims = current_claims()
    if not claims:
        return None

    # Prefer nice, human-readable identifiers
    for key in ("preferred_username", "upn", "email", "unique_name", "name"):
        val = claims.get(key)
        if val:
            return _sanitize(val)

    # Fallback to original tid:oid scheme if nothing else is available
    tid = claims.get("tid")
    oid = claims.get("oid")
    if tid and oid:
        return f"{_sanitize(tid)}:{_sanitize(oid)}"

    return None


def get_blob_path(user_key: str, session_id: str, filename: str = "") -> str:
    """
    Build a blob path under the new layout:

      user-sessions/<user_key>/<session_id>/<filename>

    - <user_key> is current_user_id() (usually email or name, sanitized)
    - <session_id> is the existing session_... UUID
    - <filename> is optional (summaries.json, chunks.json, uploads/..., etc.)
    """
    safe_user = _sanitize(user_key)
    safe_sess = _sanitize(session_id)

    base = f"user-sessions/{safe_user}/{safe_sess}"
    if filename:
        return f"{base}/{filename}"
    return base


def upload_blob(data: bytes, path: str):
    blob_service.get_blob_client(BLOB_CONTAINER_NAME, path).upload_blob(
        data,
        overwrite=True,
    )


def download_blob(path: str) -> bytes:
    return (
        blob_service.get_blob_client(BLOB_CONTAINER_NAME, path)
        .download_blob()
        .readall()
    )


def list_user_sessions(user_key: str) -> list[str]:
    """
    Return this user's session IDs ordered by most recently used.

    New layout:
      user-sessions/<user_key>/<session_id>/...

    We approximate "most recently used" by the latest blob
    last_modified timestamp for any blob under that session folder.
    """
    safe_user = _sanitize(user_key)
    prefix = f"user-sessions/{safe_user}/"

    sessions_latest: dict[str, dt.datetime] = {}
    sessions_raw: set[str] = set()

    for blob in container.list_blobs(name_starts_with=prefix):
        # blob.name looks like "user-sessions/<user>/<session_id>/something"
        # Strip off the "user-sessions/<user>/" prefix
        if not blob.name.startswith(prefix):
            continue

        remainder = blob.name[len(prefix):]
        if "/" not in remainder:
            continue

        sess = remainder.split("/")[0]
        if not sess:
            continue

        sessions_raw.add(sess)

        ts = getattr(blob, "last_modified", None)
        if ts is None:
            continue

        prev = sessions_latest.get(sess)
        if (prev is None) or (ts > prev):
            sessions_latest[sess] = ts

    # If we have timestamps, sort by last_modified (descending = most recent first)
    if sessions_latest:
        return [
            sid
            for sid, _ in sorted(
                sessions_latest.items(),
                key=lambda kv: kv[1],
                reverse=True,
            )
        ]

    # Fallback: lexicographic, reverse (older behavior)
    return sorted(sessions_raw, reverse=True)


# ---------- display-name helpers ---------- #
def meta_path(user_key: str, session_id: str) -> str:
    return get_blob_path(user_key, session_id, "meta.json")


def set_display_name(user_key: str, session_id: str, name: str):
    upload_blob(
        json.dumps({"display_name": name}).encode("utf-8"),
        meta_path(user_key, session_id),
    )


def get_display_name(user_key: str, session_id: str) -> str:
    try:
        raw = download_blob(meta_path(user_key, session_id)).decode(
            "utf-8",
            errors="ignore",
        )
        return json.loads(raw).get("display_name", session_id)
    except Exception:
        return session_id

# ------------------------------------------------
# 6.  AI helpers
# ------------------------------------------------

def chunk_text(txt: str, size: int = 2000, overlap: int = 150):
    """Split long strings into overlapping chunks for embedding."""
    if not txt:
        return []
    if size <= 0:
        size = 2000
    if overlap < 0:
        overlap = 0
    step = max(1, size - overlap)
    return [txt[i : i + size] for i in range(0, len(txt), step)]


def batch_embed(texts: list[str], batch: int = 16):
    """Embed a list of strings in batches to stay under rate limits."""
    vecs: list[list[float]] = []
    if not texts:
        return vecs

    for i in range(0, len(texts), batch):
        part = texts[i : i + batch]
        resp = llm_client.embeddings.create(input=part, model=EMBED_DEPLOYMENT)
        vecs.extend(r.embedding for r in resp.data)

    return vecs


# ---------- generic extractors (no external tools) ---------- #
def extract_text_from_file(filename: str, data: bytes) -> str:
    """
    Return raw text from common Office / text formats
    (.txt, .docx, .xlsx, .csv, .pptx). Falls back to '' on failure.
    """
    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == ".txt":
            try:
                return data.decode("utf-8")
            except UnicodeDecodeError:
                return data.decode("latin-1", errors="ignore")

        elif ext == ".docx":
            doc = DocxDocument(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs)

        elif ext == ".xlsx":
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    rows.append("\t".join("" if c is None else str(c) for c in row))
            return "\n".join(rows)

        elif ext == ".csv":
            s = io.StringIO(data.decode("utf-8", errors="ignore"))
            out = io.StringIO()
            for r in csv.reader(s):
                out.write("\t".join(r) + "\n")
            return out.getvalue()

        elif ext == ".pptx":
            prs = Presentation(io.BytesIO(data))
            txt = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        txt.append(shape.text)
            return "\n".join(txt)

    except Exception as e:
        logger.exception(e)

    return ""  # unsupported / failed


def _llm_extract_focus(txt: str, focus_instructions: str, *, max_tokens: int = 600) -> str:
    """
    Focused extraction pass used to prevent losing key details during map-reduce summarization.
    Returns a compact, structured extraction strictly from the provided text.
    """
    if not txt or not txt.strip():
        return ""

    try:
        r = llm_client.chat.completions.create(
            model=SUMMARY_DEPLOYMENT,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You extract information strictly from the provided text. "
                        "Do not invent or infer missing values. "
                        "If an item is not present, omit it."
                    ),
                },
                {
                    "role": "user",
                    "content": f"{focus_instructions}\n\nTEXT:\n{txt}",
                },
            ],
            temperature=0.0,
            max_tokens=max_tokens,
        )
        return (r.choices[0].message.content or "").strip()
    except Exception as e:
        logger.exception("Focused extraction failed: %s", e)
        return ""


def _extract_numbers_and_dates(txt: str) -> str:
    """
    Extract a compact list of financials + wages + benefits + reserves + key numeric values + dates.
    This is used as an additional input to the final synthesis so numbers don't get dropped.
    """
    # Chunk a bit smaller so we don't miss scattered wage lines across the file
    chunks = chunk_text(txt, size=8000, overlap=250)
    if not chunks:
        return ""

    focus = (
        "Extract ONLY numeric/financial and timeline facts. "
        "Return bullet points grouped under these headings (omit headings that have no items):\n"
        "**Wages / Earnings** (AWW, hourly rate, salary, wage statements, gross/net if stated)\n"
        "**Benefits / Indemnity** (TTD/TPD/PPD, weeks/dates paid, rates)\n"
        "**Medical / Bills** (amounts, dates, providers if stated)\n"
        "**Reserves / Payments / Recoveries** (reserves, payments, subrogation, liens, legal costs; include amounts if stated)\n"
        "**Settlement / Demands / Offers** (amounts and dates if stated)\n"
        "**Key Dates** (injury date, treatment dates, filing/hearing dates, RTW/termination dates if stated)\n\n"
        "Rules:\n"
        "- Use exact amounts/dates as written.\n"
        "- If an amount/date is referenced but not stated, do NOT guess; omit it.\n"
        "- Keep it concise but include all amounts/dates you find."
    )

    max_workers = int(os.getenv("SUMMARY_FOCUS_WORKERS", "6"))
    per_chunk_tokens = int(os.getenv("SUMMARY_FOCUS_MAX_TOKENS", "450"))

    extracted: list[str] = []
    try:
        with ThreadPoolExecutor(max_workers=max_workers) as ex:
            for out in ex.map(lambda c: _llm_extract_focus(c, focus, max_tokens=per_chunk_tokens), chunks):
                if out:
                    extracted.append(out)
    except Exception as e:
        logger.exception("Numbers/dates extraction failed: %s", e)
        return ""

    if not extracted:
        return ""

    # Light synthesis of the extracted bullets into a single consolidated block
    consolidate = (
        "Consolidate the following extracted bullets into a single clean set of bullets. "
        "Remove duplicates. Keep exact amounts and dates. Do not add anything new.\n\n"
        f"{'\n\n---\n\n'.join(extracted)}"
    )
    return _llm_extract_focus(consolidate, "Return the consolidated bullets only.", max_tokens=700)


def get_summary(txt: str, max_len: int = 12000) -> str:
    """
    Summarise text, automatically chunking very large documents.

    Improvements:
    - Applies SYSTEM_PROMPT_SUMMARY to chunk summaries (so role rules remain consistent).
    - Adds a focused extraction pass for numbers/dates so financial details don't get dropped.
    """
    if not txt or not txt.strip():
        return "Document contains no readable text."

    # ---------- Small document ----------
    if len(txt) < max_len:
        try:
            r = llm_client.chat.completions.create(
                model=SUMMARY_DEPLOYMENT,
                messages=[
                    {"role": "system", "content": SYSTEM_PROMPT_SUMMARY},
                    {"role": "user", "content": txt},
                ],
                temperature=0.2,
                max_tokens=1200,
            )
            return r.choices[0].message.content
        except Exception:
            return "Error: summary failed."

    # ---------- Large document (map-reduce + focused extraction) ----------
    chunks = chunk_text(txt, 10000, 200)
    if not chunks:
        return "Document unreadable or empty."

    # Focused extraction to protect numbers/dates from being dropped
    numbers_and_dates = _extract_numbers_and_dates(txt)

    partials: list[str] = []
    max_workers = int(os.getenv("SUMMARY_MAP_WORKERS", "5"))
    per_chunk_tokens = int(os.getenv("SUMMARY_CHUNK_MAX_TOKENS", "450"))

    try:
        with ThreadPoolExecutor(max_workers=max_workers) as ex:
            for out in ex.map(
                lambda c: llm_client.chat.completions.create(
                    model=SUMMARY_DEPLOYMENT,
                    messages=[
                        {"role": "system", "content": SYSTEM_PROMPT_SUMMARY},
                        {"role": "user", "content": c},
                    ],
                    temperature=0.2,
                    max_tokens=per_chunk_tokens,
                ).choices[0].message.content,
                chunks,
            ):
                if out:
                    partials.append(out)
    except Exception:
        logger.exception("Chunk summarization failed.")
        return "Error: summarisation chunks failed."

    if not partials:
        return "Error: summarisation chunks failed."

    combined = "\n\n---\n\n".join(partials)

    # Final synthesis: give the model BOTH the narrative summaries and the extracted numbers/dates
    try:
        synthesis_payload = (
            "Synthesize a single final summary strictly following the system instructions.\n\n"
            "PARTIAL SUMMARIES (narrative):\n"
            f"{combined}\n\n"
        )
        if numbers_and_dates.strip():
            synthesis_payload += (
                "EXTRACTED NUMBERS & KEY DATES (must be reflected in Numbers/Financials and Timeline if applicable):\n"
                f"{numbers_and_dates}\n\n"
            )

        r = llm_client.chat.completions.create(
            model=SUMMARY_DEPLOYMENT,
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT_SUMMARY},
                {"role": "user", "content": synthesis_payload},
            ],
            temperature=0.2,
            max_tokens=1600,
        )
        return r.choices[0].message.content
    except Exception:
        logger.exception("Final synthesis failed.")
        return "Error: final synthesis failed."


def _analyze_pdf_bytes_with_di(
    filename: str, pdf_bytes: bytes, page_offset: int = 0
) -> tuple[str, list[dict]]:
    """
    Call Document Intelligence on ONE PDF chunk (already split by pages).

    Returns:
      - full markdown content for this chunk
      - list of per-page/text chunks for FAISS with correct page numbers
    """
    body = {"base64Source": base64.b64encode(pdf_bytes).decode("ascii")}

    poller = doc_client.begin_analyze_document(
        "prebuilt-read",
        body,
        output_content_format=DocumentContentFormat.MARKDOWN,
    )
    result = poller.result()

    text = result.content or ""
    chunks: list[dict] = []

    if getattr(result, "pages", None):
        for p in result.pages:
            lines = getattr(p, "lines", None) or []
            page_text = "\n".join(getattr(l, "content", "") for l in lines)

            for ch in chunk_text(page_text):
                chunks.append(
                    {
                        "source_filename": filename,
                        "page_number": page_offset + p.page_number,
                        "text": ch,
                    }
                )

    return text, chunks


def process_pdf(filename: str, data: bytes):
    """
    PDF pipeline using Azure Document Intelligence **with pre-chunking**.

    - Splits the PDF by pages into smaller PDFs so each DI call stays under limits.
    - Calls DI per chunk, concatenates markdown, summarizes combined text,
      builds FAISS chunks with correct page numbers.
    """
    try:
        reader = PdfReader(io.BytesIO(data))
        total_pages = len(reader.pages)
    except Exception as e:
        logger.exception("Failed to parse PDF %s", filename)
        return (
            {
                "source_filename": filename,
                "summary": f"Error reading PDF: {e}",
                "error": True,
            },
            [],
        )

    if total_pages == 0:
        return (
            {
                "source_filename": filename,
                "summary": "Document has no pages or is unreadable.",
                "error": True,
            },
            [],
        )

    max_pages_per_chunk = int(os.getenv("DI_MAX_PAGES_PER_CHUNK", "25"))
    if max_pages_per_chunk <= 0 or max_pages_per_chunk > 500:
        max_pages_per_chunk = 250

    logger.info(
        "Processing PDF %s with %d pages using chunks of %d pages",
        filename,
        total_pages,
        max_pages_per_chunk,
    )

    all_text_segments: list[str] = []
    all_chunks: list[dict] = []

    page_start = 0
    while page_start < total_pages:
        page_end = min(total_pages, page_start + max_pages_per_chunk)

        writer = PdfWriter()
        for i in range(page_start, page_end):
            writer.add_page(reader.pages[i])

        buf = io.BytesIO()
        writer.write(buf)
        pdf_chunk_bytes = buf.getvalue()

        logger.info(
            "Sending PDF chunk %s pages [%d, %d) to Document Intelligence (%d bytes)",
            filename,
            page_start + 1,
            page_end,
            len(pdf_chunk_bytes),
        )

        try:
            text_segment, chunk_list = _analyze_pdf_bytes_with_di(
                filename, pdf_chunk_bytes, page_offset=page_start
            )
        except HttpResponseError as e:
            logger.exception(
                "Document Intelligence HTTP error on chunk %s [%d, %d)",
                filename,
                page_start + 1,
                page_end,
            )
            all_text_segments.append(
                f"[Error from Document Intelligence on pages {page_start+1}-{page_end}: "
                f"{getattr(e, 'message', str(e))}]"
            )
            page_start = page_end
            continue
        except Exception as e:
            logger.exception(
                "Unexpected error in DI for chunk %s [%d, %d)",
                filename,
                page_start + 1,
                page_end,
            )
            all_text_segments.append(
                f"[Unexpected error on pages {page_start+1}-{page_end}: {e}]"
            )
            page_start = page_end
            continue

        if text_segment:
            all_text_segments.append(text_segment)
        if chunk_list:
            all_chunks.extend(chunk_list)

        page_start = page_end

    if not all_text_segments:
        return (
            {
                "source_filename": filename,
                "summary": "Document unreadable or all chunks failed.",
                "error": True,
            },
            [],
        )

    combined_text = "\n\n".join(all_text_segments)
    summary = get_summary(combined_text)
    card = {"source_filename": filename, "summary": summary}
    return (card, all_chunks)


def process_file(filename: str, data: bytes):
    """
    Route: PDFs via DI with page-based chunking; other files via local text extraction.
    """
    if filename.lower().endswith(".pdf"):
        return process_pdf(filename, data)

    text = extract_text_from_file(filename, data)
    if not text.strip():
        return (
            {"source_filename": filename, "summary": "Document unreadable."},
            [],
        )

    summary = get_summary(text)
    card = {"source_filename": filename, "summary": summary}
    chunks = [
        {"source_filename": filename, "page_number": 1, "text": ch}
        for ch in chunk_text(text)
    ]
    return (card, chunks)

# ------------------------------------------------
# 7.  Auth routes (AAD)
# ------------------------------------------------
@app.route("/login")
def login():
    """
    Start Azure AD auth:
    - Uses fixed AAD_REDIRECT_URI so it doesn't depend on Host headers
      or Azure front-door / custom domain behaviour.
    - Stores MSAL flow state in the session.
    """
    try:
        # For local development, build a redirect URI from the incoming request
        # (scheme + host, which includes port) so local hosts like
        # http://localhost:8080/get_token work without changing APP_BASE_URL.
        # In production we fall back to the environment-configured AAD_REDIRECT_URI.
        host = None
        try:
            host = request.host or request.environ.get("HTTP_HOST", "")
        except RuntimeError:
            host = ""

        if host and "localhost" in host:
            scheme = request.scheme or ("https" if AAD_REDIRECT_URI.startswith("https") else "http")
            redirect_uri = f"{scheme}://{host}{AAD_REDIRECT}"
        else:
            redirect_uri = AAD_REDIRECT_URI

        logger.info("Using redirect URI for auth: %s", redirect_uri)

        session["flow"] = msal_client.initiate_auth_code_flow(
            scopes=AAD_SCOPE,
            redirect_uri=redirect_uri,
        )
    except Exception as e:
        logger.exception("Failed to initiate auth code flow: %s", e)
        return jsonify(error="auth_init_failed", message="An internal error occurred during authentication."), 500

    auth_uri = session["flow"].get("auth_uri")
    if not auth_uri:
        logger.error("MSAL did not return auth_uri in flow.")
        return jsonify(error="auth_uri_missing"), 500

    return redirect(auth_uri)


@app.route(AAD_REDIRECT)
def authorized():
    """
    Azure AD redirect target:
    - Validates the auth code using MSAL.
    - On success, stores id_token_claims in session["user"].
    - On failure, returns a clear message instead of silently failing.
    """
    flow = session.get("flow")
    if not flow:
        # This is the classic symptom of session issues (different instance, bad secret key, etc.)
        logger.warning("Auth callback with no flow in session – possible session/instance mismatch.")
        return (
            "Login failed: missing auth state. Please close the tab and try again. "
            "If this keeps happening, contact support.",
            400,
        )

    try:
        result = msal_client.acquire_token_by_auth_code_flow(flow, request.args)
    except ValueError as e:
        # Typically state mismatch, reused code, wrong redirect URI, etc.
        logger.exception("MSAL auth_code_flow ValueError: %s", e)
        return (
            "Login failed: invalid auth state (ValueError). "
            "Please try again. If it persists, contact support.",
            400,
        )
    except Exception as e:
        logger.exception("MSAL auth_code_flow unexpected error: %s", e)
        return (
            "Login failed due to an unexpected error during sign-in. "
            "Please try again later.",
            500,
        )

    if "error" in result:
        logger.error(
            "MSAL error during token acquisition: %s | desc=%s",
            result.get("error"),
            result.get("error_description"),
        )
        return (
            f"Login error from identity provider: {result.get('error_description')}",
            400,
        )

    # Store full id_token claims for later; current_user_id derives UID from tid+oid.
    claims = result.get("id_token_claims")
    if not isinstance(claims, dict):
        logger.error("id_token_claims missing or invalid in MSAL result.")
        return (
            "Login failed: did not receive valid identity information. "
            "Please try again.",
            400,
        )

    session["user"] = claims
    # Clear flow so it can't be reused accidentally.
    session.pop("flow", None)

    return redirect(url_for("index_page"))


@app.route("/logout")
def logout():
    """
    Clear local session and redirect to Azure AD logout endpoint, then back to app.
    """
    session.clear()
    post_logout = url_for("index_page", _external=True)
    return redirect(
        f"{AAD_AUTHORITY}/oauth2/v2.0/logout"
        f"?post_logout_redirect_uri={post_logout}"
    )

# ------------------------------------------------
# 8.  Async jobs + API – process files (create session)
# ------------------------------------------------

# Simple job storage prefix in Blob. Kept separate from METRICS_PREFIX so
# you don't mix operational jobs with analytics. Override via env if needed.
JOBS_PREFIX = os.getenv("JOBS_PREFIX", "jobs").strip("/") or "jobs"

# Background executor for long-running document jobs.
# This is NOT tied to request lifetime, so jobs can outlive the original HTTP call.
JOB_MAX_WORKERS = int(os.getenv("JOB_MAX_WORKERS", "2"))
JOB_EXECUTOR = ThreadPoolExecutor(max_workers=JOB_MAX_WORKERS)


def _job_blob_path(user_key: str, job_id: str) -> str:
    """
    Blob path for a job record:

        jobs/<user_key>/<job_id>.json

    Use JOBS_PREFIX env var to change the root folder if desired.
    """
    safe_user = _sanitize(user_key)
    safe_job = _sanitize(job_id)
    prefix = JOBS_PREFIX or "jobs"
    return f"{prefix}/{safe_user}/{safe_job}.json"


def _write_job_status(
    user_key: str,
    job_id: str,
    *,
    session_id: str,
    status: str,
    message: str = "",
    error: str | None = None,
):
    """
    Persist a job status snapshot to Blob.

    status: "pending" | "processing" | "completed" | "failed"
    """
    now = dt.datetime.utcnow().isoformat() + "Z"
    job_doc = {
        "job_id": job_id,
        "session_id": session_id,
        "user_id": user_key,
        "status": status,
        "message": message,
        "error": error,
        "updated_utc": now,
    }
    upload_blob(
        json.dumps(job_doc).encode("utf-8"),
        _job_blob_path(user_key, job_id),
    )


def _run_session_job(
    user_key: str,
    session_id: str,
    file_data: list[tuple[str, bytes]],
    total_bytes: int,
    job_id: str,
):
    """
    Background worker that does what the old /api/process_files SSE stream did:

      - process_file(...) for each upload (with DI for PDFs)
      - write summaries.json, chunks.json, index.faiss
      - emit metrics

    Runs in a ThreadPoolExecutor, completely decoupled from HTTP request lifetime.
    """
    started_at = dt.datetime.utcnow()
    _write_job_status(
        user_key,
        job_id,
        session_id=session_id,
        status="processing",
        message="Processing documents",
    )

    summaries: dict[str, dict] = {}
    all_chunks: list[dict] = []

    try:
        # Process each file in parallel just like before.
        with ThreadPoolExecutor(max_workers=4) as ex:
            fut = {ex.submit(process_file, n, d): n for n, d in file_data}
            for future in as_completed(fut):
                name = fut[future]
                try:
                    card, ch = future.result()
                except Exception as e:
                    logger.exception("Error processing file %s in job %s", name, job_id)
                    # Record an error card so the UI has *something* to show
                    err_card = {
                        "source_filename": name,
                        "summary": f"Error processing document: {e}",
                        "error": True,
                    }
                    summaries[name] = err_card
                    continue

                summaries[card["source_filename"]] = card
                if ch:
                    all_chunks.extend(ch)

        _write_job_status(
            user_key,
            job_id,
            session_id=session_id,
            status="processing",
            message="Saving session to storage",
        )

        # Persist summaries + chunks
        upload_blob(
            json.dumps(summaries).encode("utf-8"),
            get_blob_path(user_key, session_id, "summaries.json"),
        )
        upload_blob(
            json.dumps(all_chunks).encode("utf-8"),
            get_blob_path(user_key, session_id, "chunks.json"),
        )

        # Build FAISS index if we have any chunks
        if all_chunks:
            _write_job_status(
                user_key,
                job_id,
                session_id=session_id,
                status="processing",
                message="Building search index",
            )
            vecs = np.array(
                batch_embed([c["text"] for c in all_chunks]),
                dtype="float32",
            )
            faiss.normalize_L2(vecs)
            index = faiss.IndexFlatL2(vecs.shape[1])
            index.add(vecs)
            buf = io.BytesIO()
            faiss.write_index(index, faiss.PyCallbackIOWriter(buf.write))
            upload_blob(
                buf.getvalue(),
                get_blob_path(user_key, session_id, "index.faiss"),
            )

        # Metrics: session processing finished
        duration_ms = (dt.datetime.utcnow() - started_at).total_seconds() * 1000.0
        _safe_write_metrics(
            "session_processing_completed",
            {
                "session_id": session_id,
                "documents_processed": len(summaries),
                "total_chunks": len(all_chunks),
                "duration_ms": duration_ms,
                "total_bytes": total_bytes,
                "job_id": job_id,
            },
        )

        _write_job_status(
            user_key,
            job_id,
            session_id=session_id,
            status="completed",
            message="Session ready",
            error=None,
        )

    except Exception as e:
        logger.exception("Job %s failed for session %s", job_id, session_id)
        duration_ms = (dt.datetime.utcnow() - started_at).total_seconds() * 1000.0
        _safe_write_metrics(
            "session_processing_failed",
            {
                "session_id": session_id,
                "duration_ms": duration_ms,
                "total_bytes": total_bytes,
                "job_id": job_id,
                "error": str(e),
            },
        )
        _write_job_status(
            user_key,
            job_id,
            session_id=session_id,
            status="failed",
            message="Job failed",
            error=str(e),
        )


@app.route("/api/process_files", methods=["POST"])
def process_files():
    """
    New async behavior:

      1) Validate user + files.
      2) Create a new session_id and upload raw files to Blob.
      3) Enqueue a background job to process the documents.
      4) Return JSON { job_id, session_id } immediately.

    The front-end should then poll /api/job_status/<job_id> until
    status == "completed", and finally call /api/session_data/<session_id>
    to get summaries + chat history as usual.

    This avoids long-running HTTP connections and 504s for massive PDFs.
    """
    uid = current_user_id()
    if not uid:
        return Response(
            json.dumps({"error": "auth"}),
            status=401,
            mimetype="application/json",
        )

    files = request.files.getlist("files")
    if not files:
        return Response(
            json.dumps({"error": "no files"}),
            status=400,
            mimetype="application/json",
        )

    # Per-request safety limits (override via env in Azure)
    max_files = int(os.getenv("MAX_UPLOAD_FILES", "20"))

    # Optional size limits (MB). If set to 0 or negative, the check is disabled.
    max_file_mb_env = int(os.getenv("MAX_UPLOAD_FILE_MB", "0"))
    max_total_mb_env = int(os.getenv("MAX_UPLOAD_TOTAL_MB", "0"))

    max_file_bytes = max_file_mb_env * 1024 * 1024 if max_file_mb_env > 0 else None
    max_total_bytes = max_total_mb_env * 1024 * 1024 if max_total_mb_env > 0 else None

    if len(files) > max_files:
        return Response(
            json.dumps(
                {
                    "error": "too_many_files",
                    "detail": f"Max {max_files} files per request.",
                }
            ),
            status=400,
            mimetype="application/json",
        )

    custom_name = (request.form.get("display_name") or "").strip()

    session_id = f"session_{uuid.uuid4()}"
    default_name = dt.datetime.utcnow().strftime("Session %Y-%m-%d %H:%M")
    set_display_name(uid, session_id, custom_name or default_name)
    logger.info("Created session %s for user %s", session_id, uid)

    # Metrics: session created + upload attempt
    _safe_write_metrics(
        "session_created",
        {
            "session_id": session_id,
            "files_count": len(files),
            "filenames": [f.filename for f in files],
        },
    )

    # Read all files into memory and enforce size limits up-front.
    total_bytes = 0
    tmp_files: list[tuple[str, bytes]] = []

    for f in files:
        data = f.read() or b""
        size = len(data)

        if size == 0:
            return Response(
                json.dumps(
                    {
                        "error": "empty_file",
                        "detail": f"File '{f.filename}' is empty.",
                    }
                ),
                status=400,
                mimetype="application/json",
            )

        if max_file_bytes is not None and size > max_file_bytes:
            return Response(
                json.dumps(
                    {
                        "error": "file_too_large",
                        "detail": (
                            f"File '{f.filename}' exceeds {max_file_mb_env} MB limit. "
                            "Increase MAX_UPLOAD_FILE_MB or set it to 0 to disable this check."
                        ),
                    }
                ),
                status=413,
                mimetype="application/json",
            )

        total_bytes += size
        if max_total_bytes is not None and total_bytes > max_total_bytes:
            return Response(
                json.dumps(
                    {
                        "error": "request_too_large",
                        "detail": (
                            f"Combined upload exceeds {max_total_mb_env} MB limit. "
                            "Increase MAX_UPLOAD_TOTAL_MB or set it to 0 to disable this check."
                        ),
                    }
                ),
                status=413,
                mimetype="application/json",
            )

        tmp_files.append((f.filename, data))

    if not tmp_files:
        return Response(
            json.dumps({"error": "no_valid_files"}),
            status=400,
            mimetype="application/json",
        )

    # Persist raw uploads to Blob so they're available even if the process restarts.
    file_data: list[tuple[str, bytes]] = []
    for name, data in tmp_files:
        upload_blob(data, get_blob_path(uid, session_id, f"uploads/{name}"))
        logger.info("Uploaded file for async processing: %s (%d bytes)", name, len(data))
        file_data.append((name, data))

    # Create a job and kick off background processing
    job_id = f"job_{uuid.uuid4()}"
    logger.info("Enqueuing job %s for session %s (user %s)", job_id, session_id, uid)

    _write_job_status(
        uid,
        job_id,
        session_id=session_id,
        status="pending",
        message="Job queued",
    )

    # Run the heavy processing in the background.
    JOB_EXECUTOR.submit(
        _run_session_job,
        uid,
        session_id,
        file_data,
        total_bytes,
        job_id,
    )

    # Front-end now polls /api/job_status/<job_id> and later /api/session_data/<session_id>.
    return Response(
        json.dumps(
            {
                "job_id": job_id,
                "session_id": session_id,
                "status": "queued",
            }
        ),
        status=202,
        mimetype="application/json",
    )


@app.route("/api/job_status/<job_id>", methods=["GET"])
def job_status(job_id: str):
    """
    Check the status of an async session-processing job.

    Returns JSON like:
      {
        "job_id": "...",
        "session_id": "...",
        "user_id": "...",
        "status": "pending" | "processing" | "completed" | "failed",
        "message": "...",
        "error": null | "error text",
        "updated_utc": "..."
      }
    """
    uid = current_user_id()
    if not uid:
        return jsonify(error="auth"), 401

    path = _job_blob_path(uid, job_id)
    try:
        raw = download_blob(path).decode("utf-8", errors="ignore")
        doc = json.loads(raw)
        return jsonify(doc)
    except ResourceNotFoundError:
        return jsonify(error="unknown_job"), 404
    except Exception as e:
        logger.exception("Failed to read job status for %s", job_id)
        return jsonify(error="internal_error", detail="An internal error has occurred."), 500

# ------------------------------------------------
# 9.  API – ask question (with memory, streamed)
# ------------------------------------------------
@app.route("/api/ask_question", methods=["POST"])
def ask():
    uid = current_user_id()
    if not uid:
        return Response(json.dumps({"error":"auth"}), status=401, mimetype="application/json")

    data = request.get_json() or {}
    sess = (data.get("session_id") or "").strip()
    q    = (data.get("question")   or "").strip()

    if not sess:
        return Response(json.dumps({"error":"session_id required"}), status=400, mimetype="application/json")
    if not q:
        return Response(json.dumps({"error":"question required"}), status=400, mimetype="application/json")

    # Hard gate: only allow if session belongs to this uid
    user_sess = set(list_user_sessions(uid))
    if sess not in user_sess:
        return Response(json.dumps({"error":"unknown session"}), status=404, mimetype="application/json")

    started_at = dt.datetime.utcnow()
    question_id = uuid.uuid4().hex

    # Metrics: question received
    _safe_write_metrics(
        "qa_question_received",
        {
            "session_id": sess,
            "question_id": question_id,
            "question_length": len(q),
        },
    )

    # Retrieval knobs (tune via Azure App Settings)
    QA_FAISS_K = int(os.getenv("QA_FAISS_K", "12"))
    QA_MAX_CONTEXT_CHARS = int(os.getenv("QA_MAX_CONTEXT_CHARS", "18000"))

    try:
        idx_bytes = download_blob(get_blob_path(uid, sess, "index.faiss"))
        index = faiss.read_index(faiss.PyCallbackIOReader(io.BytesIO(idx_bytes).read))
        chunks = json.loads(
            download_blob(get_blob_path(uid, sess, "chunks.json")).decode("utf-8", errors="ignore")
        )

        if not chunks:
            doc_context = ""
            hits = []
        else:
            q_vec = np.array(
                [llm_client.embeddings.create(input=[q], model=EMBED_DEPLOYMENT).data[0].embedding],
                dtype="float32",
            )
            faiss.normalize_L2(q_vec)

            k = min(QA_FAISS_K, len(chunks))
            _, idx = index.search(q_vec, k=k)

            # Build context with dedupe + a hard size budget
            seen_keys: set[tuple[str, int]] = set()
            hits: list[dict] = []
            parts: list[str] = []
            total_chars = 0

            for i in idx[0]:
                if i < 0 or i >= len(chunks):
                    continue

                h = chunks[i]
                key = (h.get("source_filename", ""), int(h.get("page_number", 0) or 0))

                # Dedupe by (file, page)
                if key in seen_keys:
                    continue

                block = (
                    f"From {h['source_filename']} (Page {h['page_number']}):\n"
                    f"{h['text']}"
                )

                if total_chars + len(block) > QA_MAX_CONTEXT_CHARS:
                    break

                seen_keys.add(key)
                hits.append(h)
                parts.append(block)
                total_chars += len(block)

            doc_context = "\n\n---\n\n".join(parts)

        ch_path = get_blob_path(uid, sess, "chathistory.json")
        try:
            history = json.loads(download_blob(ch_path).decode("utf-8", errors="ignore"))
        except ResourceNotFoundError:
            history = []

        trimmed = history[-CHAT_HISTORY_TURNS:]

        messages = [
            {
                "role": "system",
                "content": f"{SYSTEM_PROMPT_CHAT}\n\nDocument context:\n{doc_context}",
            }
        ]
        for m in trimmed:
            role = "assistant" if m.get("role") == "bot" else "user"
            messages.append({"role": role, "content": m.get("text", "")})
        messages.append({"role": "user", "content": q})

        def stream():
            collected = ""
            sources = hits
            try:
                resp = llm_client.chat.completions.create(
                    model=CHAT_DEPLOYMENT,
                    messages=messages,
                    temperature=0.1,
                    max_tokens=8000,
                    stream=True,
                )
                for chunk in resp:
                    if chunk.choices and chunk.choices[0].delta.content:
                        token = chunk.choices[0].delta.content
                        collected += token
                        yield f"data: {json.dumps({'type': 'chunk', 'token': token})}\n\n"

                # save history at the end
                history.extend([
                    {"role": "user", "text": q},
                    {"role": "bot", "text": collected, "sources": sources},
                ])
                upload_blob(json.dumps(history).encode("utf-8"), ch_path)

                logger.info(
                    "Ask complete: session=%s, question_len=%d, answer_len=%d, sources=%d, context_chars=%d",
                    sess,
                    len(q),
                    len(collected),
                    len(sources or []),
                    len(doc_context or ""),
                )

                # Metrics: successful answer
                duration_ms = (dt.datetime.utcnow() - started_at).total_seconds() * 1000.0
                _safe_write_metrics(
                    "qa_answer_completed",
                    {
                        "session_id": sess,
                        "question_id": question_id,
                        "question_length": len(q),
                        "answer_length": len(collected),
                        "sources_count": len(sources or []),
                        "duration_ms": duration_ms,
                        "success": True,
                        "qa_faiss_k": QA_FAISS_K,
                        "qa_max_context_chars": QA_MAX_CONTEXT_CHARS,
                        "context_chars_used": len(doc_context or ""),
                    },
                )

                yield f"data: {json.dumps({'type': 'done', 'text': collected, 'sources': sources})}\n\n"

            except Exception as e:
                logger.exception(e)

                # Metrics: streaming / answer failure
                duration_ms = (dt.datetime.utcnow() - started_at).total_seconds() * 1000.0
                _safe_write_metrics(
                    "qa_answer_failed",
                    {
                        "session_id": sess,
                        "question_id": question_id,
                        "question_length": len(q),
                        "duration_ms": duration_ms,
                        "success": False,
                        "error": str(e),
                    },
                )

                yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"

        # IMPORTANT: keep request/session context alive for streaming
        return Response(stream_with_context(stream()), mimetype="text/event-stream")

    except Exception as e:
        logger.exception(e)

        # Metrics: internal error before streaming
        _safe_write_metrics(
            "qa_internal_error",
            {
                "session_id": sess,
                "question_id": question_id,
                "question_length": len(q),
                "error": str(e),
            },
        )

        return Response(json.dumps({"error": str(e)}), status=500, mimetype="application/json")

# ------------------------------------------------
# 10.  API – sessions list / rename
# ------------------------------------------------
@app.route("/api/sessions", methods=["GET"])
def api_sessions():
    uid = current_user_id()
    if not uid:
        return jsonify(error="auth"), 401
    try:
        return jsonify(
            [
                {"id": sid, "display_name": get_display_name(uid, sid)}
                for sid in list_user_sessions(uid)
            ]
        )
    except Exception as e:
        logger.exception(e)
        return jsonify(error=str(e)), 500

@app.route("/api/session_data/<session_id>", methods=["GET"])
def api_session_data(session_id):
    uid = current_user_id()
    if not uid:
        return jsonify(error="auth"), 401

    # Hard gate: only allow if session belongs to this uid
    if session_id not in set(list_user_sessions(uid)):
        return jsonify(error="unknown session"), 404

    try:
        sums = json.loads(
            download_blob(get_blob_path(uid, session_id, "summaries.json")).decode("utf-8", errors="ignore")
        )
        try:
            chat = json.loads(
                download_blob(get_blob_path(uid, session_id, "chathistory.json")).decode("utf-8", errors="ignore")
            )
        except ResourceNotFoundError:
            chat = []
        return jsonify(summaries=sums, chat_history=chat)
    except Exception as e:
        logger.exception(e)
        return jsonify(error=str(e)), 500

@app.route("/api/rename_session", methods=["POST"])
def api_rename():
    uid = current_user_id()
    if not uid:
        return jsonify(error="auth"), 401
    data = request.get_json() or {}
    sid, new = (data.get("session_id") or "").strip(), (data.get("new_name") or "").strip()
    if not sid or not new:
        return jsonify(error="invalid"), 400

    # Hard gate: only allow if session belongs to this uid
    if sid not in set(list_user_sessions(uid)):
        return jsonify(error="unknown session"), 404

    try:
        set_display_name(uid, sid, new)
        return jsonify(status="ok")
    except Exception as e:
        logger.exception(e)
        return jsonify(error=str(e)), 500

# ------------------------------------------------
# 11.  Add files to existing session
# ------------------------------------------------
@app.route("/api/add_files", methods=["POST"])
def add_files_to_session():
    uid = current_user_id()
    if not uid:
        return Response(json.dumps({"error": "auth"}), status=401, mimetype="application/json")

    session_id = (request.form.get("session_id") or "").strip()
    if not session_id:
        return Response(
            json.dumps({"error": "session_id required"}),
            status=400,
            mimetype="application/json",
        )

    # Hard gate: only allow if session belongs to this uid
    if session_id not in set(list_user_sessions(uid)):
        return Response(
            json.dumps({"error": "unknown session"}),
            status=404,
            mimetype="application/json",
        )

    file_objs = request.files.getlist("files")
    if not file_objs:
        return Response(
            json.dumps({"error": "no files"}), status=400, mimetype="application/json"
        )

    incoming_files_count = len(file_objs)

    # Per-request safety limits for incremental uploads
    max_files = int(os.getenv("MAX_UPLOAD_FILES", "20"))

    # Optional size limits (MB). 0 or negative disables the check.
    max_file_mb_env = int(os.getenv("MAX_UPLOAD_FILE_MB", "0"))
    max_total_mb_env = int(os.getenv("MAX_UPLOAD_TOTAL_MB", "0"))

    max_file_bytes = max_file_mb_env * 1024 * 1024 if max_file_mb_env > 0 else None
    max_total_bytes = max_total_mb_env * 1024 * 1024 if max_total_mb_env > 0 else None

    if len(file_objs) > max_files:
        return Response(
            json.dumps(
                {
                    "error": "too_many_files",
                    "detail": f"Max {max_files} files per request.",
                }
            ),
            status=400,
            mimetype="application/json",
        )

    # pull existing blobs / index
    try:
        sums = json.loads(
            download_blob(get_blob_path(uid, session_id, "summaries.json")).decode(
                "utf-8", errors="ignore"
            )
        )
        chunks_all = json.loads(
            download_blob(get_blob_path(uid, session_id, "chunks.json")).decode(
                "utf-8", errors="ignore"
            )
        )
        idx_bytes = download_blob(get_blob_path(uid, session_id, "index.faiss"))
        index = faiss.read_index(faiss.PyCallbackIOReader(io.BytesIO(idx_bytes).read))
    except Exception as e:
        logger.exception("Failed to load existing session data for %s", session_id)
        return Response(
            json.dumps({
                "error": "internal_error",
                "detail": "An internal error occurred while loading session data."
            }), status=500, mimetype="application/json"
        )

    # Validate new files before streaming / embedding
    total_bytes = 0
    tmp_files: list[tuple[str, bytes]] = []

    for f in file_objs:
        data = f.read() or b""
        size = len(data)

        if size == 0:
            return Response(
                json.dumps(
                    {
                        "error": "empty_file",
                        "detail": f"File '{f.filename}' is empty.",
                    }
                ),
                status=400,
                mimetype="application/json",
            )

        if max_file_bytes is not None and size > max_file_bytes:
            return Response(
                json.dumps(
                    {
                        "error": "file_too_large",
                        "detail": (
                            f"File '{f.filename}' exceeds {max_file_mb_env} MB limit. "
                            "Increase MAX_UPLOAD_FILE_MB or set it to 0 to disable this check."
                        ),
                    }
                ),
                status=413,
                mimetype="application/json",
            )

        total_bytes += size
        if max_total_bytes is not None and total_bytes > max_total_bytes:
            return Response(
                json.dumps(
                    {
                        "error": "request_too_large",
                        "detail": (
                            f"Combined upload exceeds {max_total_mb_env} MB limit. "
                            "Increase MAX_UPLOAD_TOTAL_MB or set it to 0 to disable this check."
                        ),
                    }
                ),
                status=413,
                mimetype="application/json",
            )

        tmp_files.append((f.filename, data))

    if not tmp_files:
        return Response(
            json.dumps({"error": "no_valid_files"}), status=400, mimetype="application/json"
        )

    # Metrics: add-files validated
    _safe_write_metrics(
        "session_add_files_validated",
        {
            "session_id": session_id,
            "incoming_files": incoming_files_count,
            "accepted_files": len(tmp_files),
            "total_bytes": total_bytes,
            "filenames": [name for name, _ in tmp_files],
        },
    )

    # Persist new uploads
    file_data: list[tuple[str, bytes]] = []
    for name, data in tmp_files:
        upload_blob(data, get_blob_path(uid, session_id, f"uploads/{name}"))
        file_data.append((name, data))

    started_at = dt.datetime.utcnow()

    def gen():
        new_chunks: list[dict] = []

        with ThreadPoolExecutor(max_workers=4) as ex:
            fut = {ex.submit(process_file, name, data): name for name, data in file_data}
            for done in as_completed(fut):
                name = fut[done]
                try:
                    card, add_chunks = done.result()
                except Exception as e:
                    logger.exception("Error processing file %s for session %s", name, session_id)
                    yield f"data: {json.dumps({'type':'error','file': name, 'message': str(e)})}\n\n"
                    continue

                sums[card["source_filename"]] = card
                if add_chunks:
                    new_chunks.extend(add_chunks)
                yield f"data: {json.dumps({'type':'summary','data':card})}\n\n"

        if new_chunks:
            vecs = np.array(batch_embed([c["text"] for c in new_chunks]), dtype="float32")
            faiss.normalize_L2(vecs)
            index.add(vecs)
            chunks_all.extend(new_chunks)

            buf = io.BytesIO()
            faiss.write_index(index, faiss.PyCallbackIOWriter(buf.write))
            upload_blob(
                buf.getvalue(),
                get_blob_path(uid, session_id, "index.faiss"),
            )
            upload_blob(
                json.dumps(chunks_all).encode("utf-8"),
                get_blob_path(uid, session_id, "chunks.json"),
            )

        upload_blob(
            json.dumps(sums).encode("utf-8"),
            get_blob_path(uid, session_id, "summaries.json"),
        )

        # Metrics: add-files completed
        duration_ms = (dt.datetime.utcnow() - started_at).total_seconds() * 1000.0
        _safe_write_metrics(
            "session_add_files_completed",
            {
                "session_id": session_id,
                "new_documents": len(file_data),
                "new_chunks": len(new_chunks),
                "total_documents": len(sums),
                "total_chunks": len(chunks_all),
                "duration_ms": duration_ms,
                "total_bytes": total_bytes,
            },
        )

        yield f"data: {json.dumps({'type':'complete','session_id':session_id})}\n\n"

    return Response(stream_with_context(gen()), mimetype="text/event-stream")

# ------------------------------------------------
# 12.  UI + PDF export
# ------------------------------------------------
def _parse_markdown_blocks(text: str) -> list[dict]:
    """
    Very lightweight markdown parser for our summaries.

    We care about:
      - headings: "#", "##", "###", or a single line like "**Title**"
      - bullet lines: "- something" or "* something"
      - normal paragraphs
      - blank lines for spacing

    Returns a list of blocks:
      { "type": "header" | "bullet" | "para" | "blank", "text": "..." }
    """
    blocks: list[dict] = []

    if not text:
        return blocks

    # Normalise newlines
    lines = text.replace("\r", "").split("\n")

    for raw in lines:
        line = raw.rstrip()
        stripped = line.strip()

        if not stripped:
            blocks.append({"type": "blank"})
            continue

        # Heading styles: "# Heading", "## Heading", "### Heading"
        if stripped.startswith(("# ", "## ", "### ")):
            heading = stripped.lstrip("#").strip()
            heading = heading.strip("* ").strip()
            if heading:
                blocks.append({"type": "header", "text": heading})
            continue

        # Heading style: "**Heading**" on its own line
        if stripped.startswith("**") and stripped.endswith("**") and len(stripped) <= 120:
            heading = stripped.strip("* ").strip()
            if heading:
                blocks.append({"type": "header", "text": heading})
            continue

        # Bullet lines: "- something" or "* something"
        if stripped.startswith("- "):
            content = stripped[2:].strip()
            content = content.replace("**", "")  # remove inline bold markers
            if content:
                blocks.append({"type": "bullet", "text": content})
            continue

        if stripped.startswith("* "):
            content = stripped[2:].strip()
            content = content.replace("**", "")
            if content:
                blocks.append({"type": "bullet", "text": content})
            continue

        # Normal paragraph; strip bold markers but keep text
        para = stripped.replace("**", "")
        blocks.append({"type": "para", "text": para})

    return blocks


def _build_pdf_from_sections(sections: list[tuple[str, str]]) -> bytes:
    """
    Build a simple, robust PDF from a list of (title, markdown_text) sections.

    We do a small markdown pass (headings + bullets) and then layout with reportlab.
    """
    from reportlab.lib.pagesizes import letter  # type: ignore
    from reportlab.pdfgen import canvas         # type: ignore
    from reportlab.lib.units import inch        # type: ignore

    import textwrap
    import io as _io

    buf = _io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    width, height = letter

    margin = 0.75 * inch
    max_width = width - 2 * margin  # kept for future if needed
    # Character-based wrapping approximation
    wrap_width = 100
    bullet_indent = 18  # points to indent bullet text

    def ensure_space(current_y: float, needed: float = 40.0) -> float:
        """If not enough vertical space on the page, start a new page."""
        if current_y - needed < margin:
            c.showPage()
            # Reset default font after page break
            c.setFont("Helvetica", 11)
            return height - margin
        return current_y

    for idx, (section_title, body_md) in enumerate(sections):
        # New page for each section except the very first
        if idx > 0:
            c.showPage()

        # Section heading at the top
        c.setFont("Helvetica-Bold", 14)
        y = height - margin
        title_text = (section_title or "Summary")[:200]
        c.drawString(margin, y, title_text)
        y -= 20

        # Parse markdown-ish content into blocks
        blocks = _parse_markdown_blocks(body_md or "")

        # Content font
        c.setFont("Helvetica", 11)

        for block in blocks:
            b_type = block.get("type")
            text = block.get("text", "")

            if b_type == "blank":
                # Just a little vertical gap
                y = ensure_space(y, needed=14)
                y -= 8
                continue

            if b_type == "header":
                y = ensure_space(y, needed=28)
                c.setFont("Helvetica-Bold", 12)
                # Wrap header if long
                header_lines = textwrap.wrap(text, width=wrap_width)
                for line in header_lines:
                    c.drawString(margin, y, line)
                    y -= 14
                y -= 6  # extra gap after a header
                c.setFont("Helvetica", 11)
                continue

            if b_type == "bullet":
                y = ensure_space(y, needed=20)
                # Wrap bullet text; slightly narrower since we prepend "• "
                bullet_lines = textwrap.wrap(text, width=wrap_width - 4)
                for i, line in enumerate(bullet_lines):
                    if y < margin:
                        c.showPage()
                        c.setFont("Helvetica", 11)
                        y = height - margin
                    if i == 0:
                        c.drawString(margin, y, u"• " + line)
                    else:
                        c.drawString(margin + bullet_indent, y, line)
                    y -= 14
                y -= 2
                continue

            # Normal paragraph
            if b_type == "para":
                y = ensure_space(y, needed=24)
                para_lines = textwrap.wrap(text, width=wrap_width)
                for line in para_lines:
                    if y < margin:
                        c.showPage()
                        c.setFont("Helvetica", 11)
                        y = height - margin
                    c.drawString(margin, y, line)
                    y -= 14
                y -= 6
                continue

        # If there were *no* blocks at all, at least say something
        if not blocks:
            y = ensure_space(y, needed=24)
            c.drawString(
                margin,
                y,
                "No summary text is available for this document.",
            )

    c.save()
    pdf_bytes = buf.getvalue()
    buf.close()
    return pdf_bytes


@app.route("/api/export_summary_pdf", methods=["POST"])
def api_export_summary_pdf():
    """
    Export a single summary (one document) as a PDF.

    Called by the front-end via a hidden POST form:
      - session_id
      - source_filename
    """
    uid = current_user_id()
    if not uid:
        return Response(json.dumps({"error": "auth"}), status=401, mimetype="application/json")

    session_id = (request.form.get("session_id") or "").strip()
    source_filename = (request.form.get("source_filename") or "").strip()

    if not session_id or not source_filename:
        return Response(
            json.dumps({"error": "missing_params"}),
            status=400,
            mimetype="application/json",
        )

    # Hard gate: only allow if session belongs to this uid
    if session_id not in set(list_user_sessions(uid)):
        return Response(
            json.dumps({"error": "unknown session"}),
            status=404,
            mimetype="application/json",
        )

    try:
        raw = download_blob(get_blob_path(uid, session_id, "summaries.json")).decode(
            "utf-8", errors="ignore"
        )
        summaries = json.loads(raw)
    except Exception as e:
        logger.exception("Failed to load summaries for export (session=%s)", session_id)
        return Response(
            json.dumps({"error": "load_failed"}),
            status=500,
            mimetype="application/json",
        )

    # summaries is a dict keyed by source_filename -> {"summary": ...}
    entry = summaries.get(source_filename)
    if not entry:
        # As a fallback, try matching by case-folded key
        lowered = {k.lower(): v for k, v in summaries.items()}
        entry = lowered.get(source_filename.lower())

    if not entry:
        return Response(
            json.dumps({"error": "not_found_for_source", "source_filename": source_filename}),
            status=404,
            mimetype="application/json",
        )

    summary_text = entry.get("summary") or ""

    # Use app branding for title & filename
    base_name = os.path.splitext(source_filename)[0] or source_filename or "document"
    display_title = f"Intelligent Document Assistant - {base_name}"
    sections = [(display_title, summary_text)]
    pdf_bytes = _build_pdf_from_sections(sections)

    safe_base = _sanitize(base_name or "summary") or "summary"
    download_name = f"Intelligent Document Assistant - {safe_base}.pdf"
    headers = {
        "Content-Type": "application/pdf",
        "Content-Disposition": f'attachment; filename="{download_name}"',
    }
    return Response(pdf_bytes, headers=headers)


@app.route("/api/export_all_summaries_pdf", methods=["POST"])
def api_export_all_summaries_pdf():
    """
    Export all summaries in a session into a single multi-section PDF.

    Called by the front-end via a hidden POST form:
      - session_id
    """
    uid = current_user_id()
    if not uid:
        return Response(json.dumps({"error": "auth"}), status=401, mimetype="application/json")

    session_id = (request.form.get("session_id") or "").strip()
    if not session_id:
        return Response(
            json.dumps({"error": "missing_session_id"}),
            status=400,
            mimetype="application/json",
        )

    # Hard gate: only allow if session belongs to this uid
    if session_id not in set(list_user_sessions(uid)):
        return Response(
            json.dumps({"error": "unknown session"}),
            status=404,
            mimetype="application/json",
        )

    try:
        raw = download_blob(get_blob_path(uid, session_id, "summaries.json")).decode(
            "utf-8", errors="ignore"
        )
        summaries = json.loads(raw)
    except Exception as e:
        logger.exception("Failed to load summaries for ALL export (session=%s)", session_id)
        return Response(
            json.dumps({"error": "load_failed"}),
            status=500,
            mimetype="application/json",
        )

    if not isinstance(summaries, dict) or not summaries:
        return Response(
            json.dumps({"error": "no_summaries"}),
            status=400,
            mimetype="application/json",
        )

    # Build sections in a stable order
    sections: list[tuple[str, str]] = []
    for src in sorted(summaries.keys()):
        entry = summaries[src] or {}
        title = f"Intelligent Document Assistant - {src}"
        sections.append((title, entry.get("summary") or ""))

    pdf_bytes = _build_pdf_from_sections(sections)

    safe_session = _sanitize(session_id or "session") or "session"
    filename = f"{safe_session}_summaries.pdf"
    headers = {
        "Content-Type": "application/pdf",
        "Content-Disposition": f'attachment; filename="{filename}"',
    }
    return Response(pdf_bytes, headers=headers)


@app.route("/")
def index_page():
    return render_template("index.html", user=session.get("user"))

# ------------------------------------------------
# 13.  QA Logging and Metrics
# ------------------------------------------------

# Top-level folder for analytics/metrics inside the blob container, e.g.:
#   analytics/events/2025-12-03/qa_answer_completed_<uuid>.json
#
# If you want a different root, set METRICS_PREFIX in env, e.g.:
#   METRICS_PREFIX=analytics
METRICS_PREFIX = os.getenv("METRICS_PREFIX", "analytics").strip("/") or "analytics"


def _current_user_metrics_info() -> dict:
    """
    Lightweight snapshot of the current user for metrics.

    This does NOT change auth; it's only for analytics/diagnostics.
    Avoid storing raw identity attributes in metrics events.
    """
    claims = current_claims() or {}
    uid = current_user_id()

    return {
        "authenticated": bool(uid or claims),
    }


def _safe_write_metrics(kind: str, payload: dict):
    """
    Best-effort metrics writer.

    Writes one JSON file per event into:
        <METRICS_PREFIX>/events/YYYY-MM-DD/<kind>_<uuid>.json

    - METRICS_PREFIX defaults to "analytics" so you get:
        analytics/events/...
    - Never raises; failures are logged as warnings only.
    """
    try:
        now = dt.datetime.utcnow()
        event_id = uuid.uuid4().hex
        date_str = now.strftime("%Y-%m-%d")

        event = {
            "kind": kind,
            "event_id": event_id,
            "timestamp_utc": now.isoformat() + "Z",
            "user": _current_user_metrics_info(),
            "payload": payload or {},
        }

        prefix = METRICS_PREFIX or "analytics"
        path = f"{prefix}/events/{date_str}/{kind}_{event_id}.json"

        upload_blob(json.dumps(event).encode("utf-8"), path)

    except Exception as e:
        # Metrics must never break prod behaviour
        logger.warning("Metrics write failed for %s: %s", kind, e)

# ------------------------------------------------
# 14.  Run
# ------------------------------------------------
if __name__ == "__main__":
    from waitress import serve

    port = int(os.getenv("PORT", "80"))
    threads = int(os.getenv("WAITRESS_THREADS", "8"))
    connection_limit = int(os.getenv("WAITRESS_CONNECTION_LIMIT", "100"))
    max_body_mb = int(os.getenv("WAITRESS_MAX_REQUEST_BODY_MB", "4096"))

    logger.info(
        "Starting Waitress on http://0.0.0.0:%s (threads=%s, connection_limit=%s, max_body=%sMB)",
        port,
        threads,
        connection_limit,
        max_body_mb,
    )

    serve(
        app,
        host="0.0.0.0",
        port=port,
        threads=threads,
        connection_limit=connection_limit,
        max_request_body_size=max_body_mb * 1024 * 1024,
    )